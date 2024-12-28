import os
import platform
from pptx import Presentation
from pdf2image import convert_from_path
from pptx.util import Inches
from werkzeug.utils import secure_filename
from PyPDF2 import PdfMerger
from flask import Flask, request, render_template, send_file
import io
import subprocess
import tempfile

app = Flask(__name__)

# Allowed file extensions
ALLOWED_EXTENSIONS = {'pdf', 'pptx'}

# Ensure the temp directory exists
TEMP_DIR = os.path.join(os.getcwd(), "temp")
os.makedirs(TEMP_DIR, exist_ok=True)

# Function to check allowed file types
def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

@app.route('/')
def home():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return 'No file part'
    file = request.files['file']
    if file.filename == '':
        return 'No selected file'
    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        return f"File uploaded successfully: {filename}"
    return 'Invalid file type'

# Helper function to convert PPT to PDF using LibreOffice (in-memory)
def convert_ppt_to_pdf(ppt_file):
    try:
        # Check if the app is running locally or on Render (cloud)
        if 'RENDER' in os.environ:  # Check if running on Render
            libreoffice_path = '/usr/bin/libreoffice'  # Path for Render (Linux environment)
        elif platform.system() == 'Windows':
            libreoffice_path = r'C:\Program Files\LibreOffice\program\soffice.exe'  # Path for Windows
        elif platform.system() == 'Linux':
            libreoffice_path = '/usr/bin/libreoffice'  # Path for Linux
        elif platform.system() == 'Darwin':  # macOS
            libreoffice_path = '/Applications/LibreOffice.app/Contents/MacOS/soffice'  # Path for macOS
        else:
            raise RuntimeError("Unsupported platform for LibreOffice")

        # Create a temporary file to save the in-memory PPT file
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_ppt_file:
            ppt_file.seek(0)  # Reset the pointer to the start of the file
            temp_ppt_file.write(ppt_file.read())
            temp_ppt_file.close()

            # Temporary output PDF file path
            temp_pdf_file = temp_ppt_file.name.replace('.pptx', '.pdf')

            # Run the LibreOffice conversion command
            command = [
                libreoffice_path,  # Full path to LibreOffice executable
                '--headless',       # Run in headless mode (no GUI)
                '--convert-to', 'pdf',  # Specify PDF format
                '--outdir', os.path.dirname(temp_pdf_file),  # Output directory
                temp_ppt_file.name  # Input PPT file (from temporary file)
            ]

            subprocess.run(command, check=True)

            # Read the generated PDF into memory
            with open(temp_pdf_file, 'rb') as f:
                pdf_output = io.BytesIO(f.read())
            pdf_output.seek(0)

            # Clean up the temporary files
            os.remove(temp_ppt_file.name)
            os.remove(temp_pdf_file)

        return pdf_output

    except Exception as e:
        raise RuntimeError(f"Error converting PPT to PDF: {e}")

# Helper function to convert PDF to PPT using LibreOffice (in-memory)
def convert_pdf_to_ppt(pdf_file):
    try:
        # Create a temporary file to save the in-memory PDF
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_pdf_file:
            temp_pdf_file.write(pdf_file.read())
            temp_pdf_path = temp_pdf_file.name

        # Convert PDF to images using pdf2image
        images = convert_from_path(temp_pdf_path, dpi=300)

        # Create a PowerPoint presentation using python-pptx
        presentation = Presentation()

        # Add each image as a new slide
        for image in images:
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])

            # Save the image to a BytesIO object
            img_stream = io.BytesIO()
            image.save(img_stream, format="PNG")
            img_stream.seek(0)

            # Add the image to the slide
            slide.shapes.add_picture(img_stream, Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))

        # Save the PowerPoint presentation to a BytesIO object
        ppt_output = io.BytesIO()
        presentation.save(ppt_output)
        ppt_output.seek(0)

        # Clean up the temporary PDF file
        os.remove(temp_pdf_path)

        return ppt_output

    except Exception as e:
        raise RuntimeError(f"Error converting PDF to PPT: {e}")



@app.route('/merge_pdfs', methods=['POST'])
def merge_pdfs():
    files = request.files.getlist('pdfs')
    if not files:
        return "No files provided for merging"
    
    merger = PdfMerger()

    # Process each uploaded file
    for file in files:
        if file and allowed_file(file.filename):
            # Read the file into memory
            file_content = file.read()
            
            # Check if the file content is empty
            if not file_content:
                return f"Error: The file {file.filename} is empty"

            # Create an in-memory BytesIO object with the file content
            file_io = io.BytesIO(file_content)
            
            # Append the in-memory file to the merger
            file_io.seek(0)  # Ensure pointer is at the start of the file
            merger.append(file_io)
        else:
            return f"Invalid file type or empty file: {file.filename}"

    # Create the merged PDF in memory
    output = io.BytesIO()
    merger.write(output)
    output.seek(0)

    # Send merged PDF as a response
    return send_file(output, as_attachment=True, download_name='merged.pdf', mimetype='application/pdf')


@app.route('/convert_ppt_to_pdf', methods=['POST'])
def convert_ppt_to_pdf_route():
    file = request.files['ppt']
    if not file or file.filename == '':
        return 'No file selected'

    filename = secure_filename(file.filename)

    # Save the uploaded PPT file into an in-memory buffer
    ppt_file = io.BytesIO(file.read())

    # Convert PPT to PDF in-memory
    try:
        pdf_output = convert_ppt_to_pdf(ppt_file)
    except Exception as e:
        return str(e)

    # Send the converted PDF file back as a response
    return send_file(pdf_output, as_attachment=True, download_name=f"{filename}.pdf", mimetype='application/pdf')


@app.route('/convert_pdf_to_ppt', methods=['POST'])
def convert_pdf_to_ppt_route():
    # Check for uploaded file
    file = request.files['pdf_to_ppt']
    if not file or file.filename == '':
        return 'No file selected'

    # Save the PDF in-memory
    pdf_file = io.BytesIO(file.read())

    try:
        # Convert PDF to PPT
        ppt_output = convert_pdf_to_ppt(pdf_file)

        # Return the PPT file in-memory without saving it to the disk
        return send_file(
            ppt_output,
            as_attachment=True,
            download_name=f"{os.path.splitext(file.filename)[0]}.pptx",
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )

    except Exception as e:
        return str(e)



@app.route('/convert_merge_ppt_to_pdf', methods=['POST'])
def convert_merge_ppt_to_pdf():
    files = request.files.getlist('ppt_multiple')
    if not files:
        return "No PPT files provided for conversion and merging"

    merger = PdfMerger()

    # Process each PPT file
    for file in files:
        filename = secure_filename(file.filename)

        # Save the uploaded PPT file into an in-memory buffer
        ppt_file = io.BytesIO(file.read())
        pdf_filename = f"{os.path.splitext(filename)[0]}.pdf"

        # Convert PPT to PDF
        try:
            pdf_output = convert_ppt_to_pdf(ppt_file)
        except Exception as e:
            return str(e)

        # Append the converted PDF to the merger
        merger.append(pdf_output)

    # Generate merged PDF in-memory
    output_pdf = io.BytesIO()
    merger.write(output_pdf)
    output_pdf.seek(0)

    # Send the merged PDF file as a response
    return send_file(output_pdf, as_attachment=True, download_name='merged_ppt_to_pdf.pdf', mimetype='application/pdf')


if __name__ == '__main__':
    app.run(debug=True)
