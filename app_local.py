from flask import Flask, request, render_template, send_file
from PyPDF2 import PdfMerger
from pptx import Presentation
from PyPDF2 import PdfReader
from reportlab.pdfgen import canvas
from io import BytesIO
import os
import aspose.slides as slides
from pptx.util import Inches, Pt
import fitz
import platform
import pythoncom
from werkzeug.utils import secure_filename
import comtypes.client

app = Flask(__name__)

UPLOAD_FOLDER = os.path.abspath('uploads')  # Ensure absolute path for uploads
ALLOWED_EXTENSIONS = {'pdf', 'pptx'}

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

# Create the uploads directory if it doesn't exist
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

@app.route('/')
def home():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return 'No file part'
    file = request.files['file']
    if file.filename == '':
        return 'No selected file'
    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(file_path)
        return f"File uploaded successfully: {filename}"
    return 'Invalid file type'

# Helper function to convert PPT to PDF (Windows only)
def convert_ppt_to_pdf(ppt_file, pdf_file):
    if platform.system() != "Windows":
        raise OSError("PowerPoint-to-PDF conversion is only supported on Windows.")
    
    if not os.path.exists(ppt_file):
        raise FileNotFoundError(f"File not found: {ppt_file}")

    try:
        pythoncom.CoInitialize()  # Initialize COM
        powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
        powerpoint.Visible = 1  # Optional: Set PowerPoint to be visible
        ppt = powerpoint.Presentations.Open(ppt_file)
        ppt.SaveAs(pdf_file, 32)  # 32 is the code for PDF format
        ppt.Close()
    except Exception as e:
        raise RuntimeError(f"Error converting {ppt_file}: {e}")
    finally:
        if 'powerpoint' in locals():
            powerpoint.Quit()
        pythoncom.CoUninitialize()  # Uninitialize COM

@app.route('/merge_pdfs', methods=['POST'])
def merge_pdfs():
    files = request.files.getlist('pdfs')
    if not files:
        return "No files provided for merging"

    merger = PdfMerger()
    for file in files:
        merger.append(file)
    output = os.path.join(app.config['UPLOAD_FOLDER'], 'merged.pdf')
    with open(output, 'wb') as f:
        merger.write(f)
    return send_file(output, as_attachment=True, download_name='merged.pdf')

@app.route('/convert_ppt_to_pdf', methods=['POST'])
def convert_ppt_to_pdf_route():
    file = request.files['ppt']
    if not file or file.filename == '':
        return 'No file selected'

    filename = secure_filename(file.filename)
    ppt_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    pdf_filename = f"{os.path.splitext(filename)[0]}.pdf"
    pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], pdf_filename)

    # Save uploaded file
    file.save(ppt_path)

    # Convert PPT to PDF
    try:
        convert_ppt_to_pdf(ppt_path, pdf_path)
    except Exception as e:
        return str(e)

    return send_file(pdf_path, as_attachment=True, download_name=pdf_filename)

@app.route('/convert_merge_ppt_to_pdf', methods=['POST'])
def convert_merge_ppt_to_pdf():
    files = request.files.getlist('ppt_multiple')
    if not files:
        return "No PPT files provided for conversion and merging"

    merger = PdfMerger()
    for file in files:
        filename = secure_filename(file.filename)
        ppt_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        pdf_filename = f"{os.path.splitext(filename)[0]}.pdf"
        pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], pdf_filename)

        # Save the uploaded PPT file
        file.save(ppt_path)

        # Convert PPT to PDF
        try:
            convert_ppt_to_pdf(ppt_path, pdf_path)
        except Exception as e:
            return str(e)

        # Append the converted PDF to the merger
        merger.append(pdf_path)

    output = os.path.join(app.config['UPLOAD_FOLDER'], 'merged_ppt_to_pdf.pdf')
    with open(output, 'wb') as f:
        merger.write(f)
    
    return send_file(output, as_attachment=True, download_name='merged_ppt_to_pdf.pdf')
@app.route('/convert_pdf_to_ppt', methods=['POST'])
def convert_pdf_to_ppt():
    file = request.files['pdf_to_ppt']  # Ensure this matches the form field name in HTML
    if not file or file.filename == '':
        return 'No file selected'

    filename = secure_filename(file.filename)
    pdf_path = os.path.join(UPLOAD_FOLDER, filename)
    ppt_filename = f"{os.path.splitext(filename)[0]}.pptx"
    ppt_path = os.path.join(UPLOAD_FOLDER, ppt_filename)

    # Save uploaded PDF file
    file.save(pdf_path)

    try:
        # Create presentation object
        with slides.Presentation() as pres:
            # Remove default slide
            pres.slides.remove_at(0)

            # Import PDF to presentation
            pres.slides.add_from_pdf(pdf_path)

            # Save the presentation as PPTX
            pres.save(ppt_path, slides.export.SaveFormat.PPTX)

        return send_file(ppt_path, as_attachment=True, download_name=ppt_filename)

    except Exception as e:
        return f"Error during conversion: {e}"
    

if __name__ == '__main__':
    app.run(debug=True)
